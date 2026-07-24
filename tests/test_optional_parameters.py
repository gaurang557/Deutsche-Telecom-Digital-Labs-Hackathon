"""An optional parameter given as null means "not using this", so it is dropped.

A live run failed on the final hop with `save_as must be a path string`. `save_as`
is optional, and a model filling in a JSON object routinely writes
`"save_as": null` for an optional field it does not want. Omitting the key is what
that means, and the executors already read an absent `save_as` as "edit the target
in place", so dropping it makes the plan runnable rather than reinterpreting it.

Two properties matter more than the fix itself and are asserted here:

  * the stripping happens BEFORE the confirmation hash is computed, so the
    parameters the user approves are the parameters that run. Doing it later would
    reopen the confirm-one-thing-do-another hole closed in an earlier round.
  * it is omission, never coercion. A string `"true"` stays the string `"true"`,
    and a REQUIRED parameter given as null is NOT quietly dropped — it earns a
    repair attempt, because a null `find` is a broken plan rather than an unused
    option.

Every fixture is built inside `tmp_path` with a throwaway home patched over
`Path.home`, so no test touches a real Desktop.
"""

# ruff: noqa: I001

import os
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.execution.hybrid import HybridExecutor, build_structured_dispatcher
from app.planning.capabilities import (
    find_non_string_path_parameter,
    find_null_required_parameter,
)
from app.planning.normalizer import build_action_plan
from app.schemas import DraftAction, DraftPlan, TaskRequest
from app.structured_actions import (
    action_identity_hash,
    is_absent_path_value,
    optional_parameters_for,
    required_parameters_for,
    strip_absent_optional_parameters,
)

from windows_agent.audit import InMemoryAuditSink


@pytest.fixture
def desktop(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    home = tmp_path / "home"
    (home / "Desktop").mkdir(parents=True)
    monkeypatch.setattr(Path, "home", classmethod(lambda cls: home))
    return home / "Desktop"


def _same(left: str | Path, right: str | Path) -> bool:
    return os.path.normcase(str(left)) == os.path.normcase(str(right))


def _deck(path: Path, wording: str) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.text_frame.paragraphs[0].text = wording
    presentation.save(str(path))
    return path


def _replace_plan(**parameters: object) -> DraftPlan:
    return DraftPlan(
        summary="I'll update the wording in that deck.",
        actions=[
            DraftAction(
                step_key="edit",
                type="presentation.replace_text",
                target="Desktop/deck.pptx",
                parameters={"find": "old wording", "replace": "new wording", **parameters},
                description="Update the deck.",
                expected_result={"replaced": True},
            )
        ],
    )


def _texts(path: Path) -> list[str]:
    """Read the deck back from DISK, never from the executor's memory."""
    presentation = Presentation(str(path))
    return [
        paragraph.text
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
    ]


# ── the stripping rule ───────────────────────────────────────────────────────
def test_a_null_optional_parameter_is_dropped() -> None:
    stripped = strip_absent_optional_parameters(
        "presentation.replace_text",
        {"find": "a", "replace": "b", "save_as": None},
    )

    assert stripped == {"find": "a", "replace": "b"}


def test_an_empty_optional_path_is_dropped() -> None:
    for blank in ("", "   ", "\t"):
        stripped = strip_absent_optional_parameters(
            "presentation.replace_text",
            {"find": "a", "replace": "b", "save_as": blank},
        )
        assert "save_as" not in stripped


def test_a_real_optional_value_is_kept() -> None:
    stripped = strip_absent_optional_parameters(
        "presentation.replace_text",
        {"find": "a", "replace": "b", "save_as": "Desktop/out.pptx", "count": 2},
    )

    assert stripped == {
        "find": "a",
        "replace": "b",
        "save_as": "Desktop/out.pptx",
        "count": 2,
    }


def test_a_null_required_parameter_is_never_dropped() -> None:
    """Dropping it would turn a clear failure into a confusing one."""
    stripped = strip_absent_optional_parameters(
        "presentation.replace_text",
        {"find": None, "replace": "b"},
    )

    assert stripped == {"find": None, "replace": "b"}


def test_an_empty_required_value_is_kept_because_it_can_be_meaningful() -> None:
    """`file.write_text` with content "" writes an empty file. That is not absence."""
    stripped = strip_absent_optional_parameters("file.write_text", {"content": ""})

    assert stripped == {"content": ""}


def test_nothing_is_coerced() -> None:
    """Strictly omission. No string becomes a bool, no string becomes a number."""
    original = {
        "find": "a",
        "replace": "b",
        "overwrite": "true",
        "count": "2",
    }

    stripped = strip_absent_optional_parameters("presentation.replace_text", original)

    assert stripped["overwrite"] == "true"
    assert isinstance(stripped["overwrite"], str)
    assert stripped["count"] == "2"
    assert isinstance(stripped["count"], str)


def test_a_legacy_action_keeps_its_parameters_untouched() -> None:
    parameters = {"anything": None, "other": ""}

    assert strip_absent_optional_parameters("open_file", parameters) == parameters


def test_optional_and_required_are_derived_from_the_registries() -> None:
    """Subtraction, not a second hand-maintained list that could drift."""
    assert "save_as" in optional_parameters_for("presentation.replace_text")
    assert "find" in required_parameters_for("presentation.replace_text")
    assert "find" not in optional_parameters_for("presentation.replace_text")
    # A legacy action has neither, which is what leaves it untouched above.
    assert optional_parameters_for("open_file") == frozenset()


# ── the plan-time repair-loop checks ─────────────────────────────────────────
def test_a_null_required_parameter_is_reported_for_repair() -> None:
    problem = find_null_required_parameter(
        "presentation.replace_text", {"find": None, "replace": "b"}
    )

    assert problem is not None
    assert "find" in problem


def test_a_null_optional_parameter_is_not_reported() -> None:
    assert (
        find_null_required_parameter(
            "presentation.replace_text",
            {"find": "a", "save_as": None},
        )
        is None
    )


def test_an_empty_required_string_is_not_called_absent() -> None:
    assert find_null_required_parameter("file.write_text", {"content": ""}) is None


@pytest.mark.parametrize("value", [0, 3, {"a": 1}, ["x"]])
def test_a_path_parameter_that_is_not_a_string_is_reported_for_repair(
    value: object,
) -> None:
    problem = find_non_string_path_parameter({"save_as": value})

    assert problem is not None
    assert repr(value) in problem


def test_a_dropped_or_real_path_parameter_is_not_reported() -> None:
    assert find_non_string_path_parameter({"save_as": None}) is None
    assert find_non_string_path_parameter({"save_as": "Desktop/out.pptx"}) is None
    assert find_non_string_path_parameter({}) is None


# ── `false` in a path parameter means "not using this" ────────────────────────
def test_false_is_absent_for_an_optional_path_parameter() -> None:
    """A path cannot be a boolean, so `save_as: false` is plainly "unused"."""
    stripped = strip_absent_optional_parameters(
        "presentation.replace_text",
        {"find": "a", "replace": "b", "save_as": False},
    )

    assert "save_as" not in stripped


def test_false_is_never_absent_for_any_other_parameter() -> None:
    """`overwrite: false` is a real, meaningful value. Workflow 1 depends on it."""
    stripped = strip_absent_optional_parameters(
        "spreadsheet.write_cell",
        {"sheet": "Revenue", "cell": "B2", "value": 27.4, "overwrite": False},
    )

    assert stripped["overwrite"] is False


def test_the_repair_check_and_the_stripping_rule_agree_about_absence() -> None:
    """They must not disagree: the check runs on the draft, BEFORE stripping.

    If the plan-time check rejected a value the stripping rule would have quietly
    removed, the repair loop would burn its budget on a plan that was about to be
    fine — which is exactly the regression that produced an opaque 422.
    """
    for value in (None, False, "", "   "):
        assert find_non_string_path_parameter({"save_as": value}) is None
        assert "save_as" not in strip_absent_optional_parameters(
            "presentation.replace_text", {"find": "a", "save_as": value}
        )


def test_zero_is_not_treated_as_an_absent_path() -> None:
    """`0 is False` is False in Python, and 0 is not a path either."""
    assert is_absent_path_value(0) is False
    assert is_absent_path_value(False) is True
    assert is_absent_path_value("Desktop/out.pptx") is False


# ── the confirmation hash covers what actually runs ──────────────────────────
def test_the_confirmation_hash_binds_the_stripped_parameters(desktop: Path) -> None:
    _deck(desktop / "deck.pptx", "old wording")

    plan = build_action_plan(
        TaskRequest(text="update the wording in Desktop/deck.pptx"),
        _replace_plan(save_as=None),
    )
    action = plan.actions[0]

    # The null never reaches the executor, and the hash is over what does.
    assert "save_as" not in action.parameters
    assert action.requires_confirmation is True
    assert action.confirmation_hash == action_identity_hash(
        "presentation.replace_text",
        action.target,
        action.parameters,
    )
    # And NOT over the unstripped set the planner sent.
    assert action.confirmation_hash != action_identity_hash(
        "presentation.replace_text",
        action.target,
        {"find": "old wording", "replace": "new wording", "save_as": None},
    )


def test_a_null_save_as_still_counts_as_an_in_place_edit(desktop: Path) -> None:
    """Stripping must not talk the plan out of needing confirmation."""
    _deck(desktop / "deck.pptx", "old wording")

    plan = build_action_plan(
        TaskRequest(text="update the wording in Desktop/deck.pptx"),
        _replace_plan(save_as=None),
    )

    assert plan.actions[0].requires_confirmation is True
    assert plan.actions[0].risk == "high"


# ── end to end ───────────────────────────────────────────────────────────────
async def test_a_null_save_as_saves_in_place_and_verifies_from_disk(
    desktop: Path,
) -> None:
    """The exact live failure, now succeeding, with the edit proved on disk."""
    deck = _deck(desktop / "deck.pptx", "old wording")

    plan = build_action_plan(
        TaskRequest(text="update the wording in Desktop/deck.pptx"),
        _replace_plan(save_as=None),
    )
    action = plan.actions[0]
    audit = InMemoryAuditSink()
    executor = HybridExecutor(dispatcher=build_structured_dispatcher(audit), audit=audit)

    response = await executor.execute_plan(
        plan,
        {action.action_id},
        approved_action_hashes={action.action_id: action.confirmation_hash or ""},
    )

    assert response.status == "completed"
    result = response.results[0]
    assert result.status == "succeeded"
    # Saved over the target rather than to some new file.
    assert _same(result.evidence["output_path"], deck)
    assert result.evidence["save_as"] is False
    # The verifier reopened the deck from disk and agreed.
    assert result.verification is not None and result.verification.passed is True
    # Independent proof, read from disk here rather than from evidence.
    assert "new wording" in _texts(deck)
    assert "old wording" not in _texts(deck)
    # No new file was created alongside it.
    assert sorted(p.name for p in desktop.iterdir()) == ["deck.pptx"]


async def test_an_empty_save_as_also_saves_in_place(desktop: Path) -> None:
    deck = _deck(desktop / "deck.pptx", "old wording")

    plan = build_action_plan(
        TaskRequest(text="update the wording in Desktop/deck.pptx"),
        _replace_plan(save_as="   "),
    )
    action = plan.actions[0]
    audit = InMemoryAuditSink()
    executor = HybridExecutor(dispatcher=build_structured_dispatcher(audit), audit=audit)

    response = await executor.execute_plan(
        plan,
        {action.action_id},
        approved_action_hashes={action.action_id: action.confirmation_hash or ""},
    )

    assert response.results[0].status == "succeeded"
    assert "new wording" in _texts(deck)
    assert sorted(p.name for p in desktop.iterdir()) == ["deck.pptx"]


async def test_a_real_save_as_still_writes_a_new_file(desktop: Path) -> None:
    """The stripping must not have broken the case save_as exists for."""
    deck = _deck(desktop / "deck.pptx", "old wording")

    plan = build_action_plan(
        TaskRequest(text="rewrite Desktop/deck.pptx into a copy"),
        _replace_plan(save_as="Desktop/out.pptx"),
    )
    action = plan.actions[0]
    audit = InMemoryAuditSink()
    executor = HybridExecutor(dispatcher=build_structured_dispatcher(audit), audit=audit)

    response = await executor.execute_plan(
        plan,
        {action.action_id},
        approved_action_hashes={action.action_id: action.confirmation_hash or ""},
    )

    assert response.results[0].status == "succeeded"
    assert "new wording" in _texts(desktop / "out.pptx")
    # The original is untouched, which is the whole point of save_as.
    assert "old wording" in _texts(deck)
