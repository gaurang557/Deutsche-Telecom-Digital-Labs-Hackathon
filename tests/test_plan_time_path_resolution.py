"""Path resolution happens at PLAN BUILD, so confirmation binds the real file.

A live run showed why this ordering matters. Two read steps were rescued onto a
file one folder down, and the third step — the one that actually modified the
deck — was handed the unresolved path and failed, because discovery only ran for
read-only actions.

The tempting fix, running discovery for modifying actions inside the executor,
would have broken the safety story instead. "Exact confirmation accepted for
step 3" is recorded BEFORE execution begins, so retargeting during execution
means the user authorised changing one file while a different file was changed.
Resolving during plan build fixes the failure without that hole: the policy
decision, the target the user is shown, and the confirmation hash all describe
the file that will really be touched.

What must NOT be resolved is equally load-bearing and asserted here: an action
that creates its target, and a `destination` / `save_as` output path, are always
taken literally, so a create can never be captured by a same-named file
elsewhere.

Every fixture is built inside `tmp_path` with a throwaway home patched over
`Path.home`, so no test touches a real Desktop. None of the folder or file names
used here is one the production code knows about.
"""

# ruff: noqa: I001

import os
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.execution.hybrid import (
    HybridExecutor,
    _discover_readable_file,
    build_structured_dispatcher,
    resolve_plan_target,
)
from app.planning.exceptions import InvalidPlannerResponseError
from app.planning.normalizer import build_action_plan
from app.schemas import DraftAction, DraftPlan, TaskRequest
from app.structured_actions import (
    CREATES_ITS_TARGET_ACTIONS,
    MODIFYING_STRUCTURED_ACTIONS,
    READ_ONLY_STRUCTURED_ACTIONS,
    REQUIRES_EXISTING_TARGET_ACTIONS,
    action_identity_hash,
)

from windows_agent.audit import InMemoryAuditSink


@pytest.fixture
def desktop(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    """A throwaway home whose Desktop is the only real allowlisted root in play."""
    home = tmp_path / "home"
    (home / "Desktop").mkdir(parents=True)
    monkeypatch.setattr(Path, "home", classmethod(lambda cls: home))
    return home / "Desktop"


def _same(left: str | Path, right: str | Path) -> bool:
    return os.path.normcase(str(left)) == os.path.normcase(str(right))


def _deck(path: Path, wording: str) -> Path:
    """A one-slide deck holding `wording` in a text box."""
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.text_frame.paragraphs[0].text = wording
    presentation.save(str(path))
    return path


def _single_step_plan(type_: str, target: str, parameters: dict) -> DraftPlan:
    return DraftPlan(
        summary="I'll take care of that for you.",
        actions=[
            DraftAction(
                step_key="step",
                type=type_,
                target=target,
                parameters=parameters,
                description="Handle the file.",
                expected_result={"contains": "the wording"},
            )
        ],
    )


# ── which actions may be resolved, derived from the canonical registry ────────
def test_every_read_requires_an_existing_target() -> None:
    assert READ_ONLY_STRUCTURED_ACTIONS <= REQUIRES_EXISTING_TARGET_ACTIONS


def test_in_place_edits_require_an_existing_target() -> None:
    # These inherently fail unless the file is already there, so resolving them
    # onto the single existing candidate cannot invent a target.
    assert "presentation.replace_text" in REQUIRES_EXISTING_TARGET_ACTIONS
    assert "document.replace_text" in REQUIRES_EXISTING_TARGET_ACTIONS


def test_actions_that_create_their_target_are_never_resolved() -> None:
    assert CREATES_ITS_TARGET_ACTIONS.isdisjoint(REQUIRES_EXISTING_TARGET_ACTIONS)
    # write_cell creates the workbook when it is missing, so it must stay out.
    assert "spreadsheet.write_cell" in CREATES_ITS_TARGET_ACTIONS
    assert "spreadsheet.write_cell" not in REQUIRES_EXISTING_TARGET_ACTIONS


def test_a_permanently_denied_action_is_never_helped_to_find_a_file() -> None:
    assert "file.delete" not in REQUIRES_EXISTING_TARGET_ACTIONS


def test_the_requires_existing_set_cannot_drift_from_the_modifying_set() -> None:
    # Derived by subtraction: every modifying action is either declared as
    # creating its target, permanently denied, or resolvable. No third list.
    unclassified = MODIFYING_STRUCTURED_ACTIONS - CREATES_ITS_TARGET_ACTIONS
    assert unclassified <= REQUIRES_EXISTING_TARGET_ACTIONS | {"file.delete"}


# ── a modifying action is resolved at plan time ───────────────────────────────
def test_a_modifying_action_target_is_resolved_during_plan_build(desktop: Path) -> None:
    actual = _deck(desktop / "handover" / "deck.pptx", "old wording")

    plan = build_action_plan(
        TaskRequest(text="update the wording in Desktop/deck.pptx"),
        _single_step_plan(
            "presentation.replace_text",
            "Desktop/deck.pptx",
            {"find": "old wording", "replace": "new wording"},
        ),
    )

    action = plan.actions[0]
    assert _same(action.target, actual)
    assert _same(action.resolved_from, desktop / "deck.pptx")


def test_a_target_that_already_exists_is_left_exactly_as_written(desktop: Path) -> None:
    _deck(desktop / "deck.pptx", "old wording")

    plan = build_action_plan(
        TaskRequest(text="update the wording in Desktop/deck.pptx"),
        _single_step_plan(
            "presentation.replace_text",
            "Desktop/deck.pptx",
            {"find": "old wording", "replace": "new wording"},
        ),
    )

    assert plan.actions[0].target == "Desktop/deck.pptx"
    assert plan.actions[0].resolved_from is None


# ── the important one: confirmation binds the RESOLVED file ───────────────────
def test_the_confirmation_hash_binds_the_resolved_path(desktop: Path) -> None:
    """The user must not approve one file and have another one changed."""
    actual = _deck(desktop / "handover" / "deck.pptx", "old wording")
    parameters = {"find": "old wording", "replace": "new wording"}

    plan = build_action_plan(
        TaskRequest(text="update the wording in Desktop/deck.pptx"),
        _single_step_plan("presentation.replace_text", "Desktop/deck.pptx", parameters),
    )

    action = plan.actions[0]
    assert action.requires_confirmation is True
    expected = action_identity_hash("presentation.replace_text", str(actual), parameters)
    assert action.confirmation_hash == expected
    # And emphatically NOT the hash of the path the planner guessed.
    stale = action_identity_hash(
        "presentation.replace_text", str(desktop / "deck.pptx"), parameters
    )
    assert action.confirmation_hash != stale


def test_the_confirmation_summary_names_the_file_that_will_change(desktop: Path) -> None:
    _deck(desktop / "handover" / "deck.pptx", "old wording")

    plan = build_action_plan(
        TaskRequest(text="update the wording in Desktop/deck.pptx"),
        DraftPlan(
            summary="I'll update that wording.",
            actions=[
                DraftAction(
                    step_key="step",
                    type="presentation.replace_text",
                    target="Desktop/deck.pptx",
                    parameters={"find": "old wording", "replace": "new wording"},
                    # No description, so the generated one is what the user sees.
                    expected_result={"replaced": True},
                )
            ],
        ),
    )

    # What the user is shown names the folder the file is really in.
    assert "handover" in plan.actions[0].description
    assert "handover" in plan.actions[0].target


# ── creates and destinations are still never redirected ───────────────────────
def test_a_create_target_is_not_redirected_at_plan_build(desktop: Path) -> None:
    (desktop / "handover").mkdir(parents=True)
    (desktop / "handover" / "book.xlsx").write_bytes(b"an existing workbook")

    target, resolved_from = resolve_plan_target(
        "spreadsheet.write_cell", "Desktop/book.xlsx"
    )

    assert target == "Desktop/book.xlsx"
    assert resolved_from is None


def test_a_save_as_output_path_is_not_redirected_at_plan_build(desktop: Path) -> None:
    _deck(desktop / "handover" / "deck.pptx", "old wording")
    _deck(desktop / "deck.pptx", "old wording")
    (desktop / "handover" / "out.pptx").write_bytes(b"decoy")

    plan = build_action_plan(
        TaskRequest(text="update the wording in Desktop/deck.pptx"),
        _single_step_plan(
            "presentation.replace_text",
            "Desktop/deck.pptx",
            {"find": "old wording", "replace": "new", "save_as": "Desktop/out.pptx"},
        ),
    )

    # resolve_plan_target only ever looks at the target; the output path is
    # untouched here and merely alias-expanded at execution time.
    assert plan.actions[0].parameters["save_as"] == "Desktop/out.pptx"


# ── ambiguity fails loudly rather than guessing ───────────────────────────────
def test_ambiguity_for_a_modifying_action_fails_with_the_candidates(desktop: Path) -> None:
    _deck(desktop / "north" / "deck.pptx", "old wording")
    _deck(desktop / "south" / "deck.pptx", "old wording")

    with pytest.raises(InvalidPlannerResponseError) as excinfo:
        build_action_plan(
            TaskRequest(text="update the wording in Desktop/deck.pptx"),
            _single_step_plan(
                "presentation.replace_text",
                "Desktop/deck.pptx",
                {"find": "old wording", "replace": "new wording"},
            ),
        )

    message = str(excinfo.value)
    assert "north" in message and "south" in message
    assert "which folder you meant" in message


# ── no double resolution ─────────────────────────────────────────────────────
def test_a_resolved_path_cannot_be_resolved_again(desktop: Path) -> None:
    actual = _deck(desktop / "handover" / "deck.pptx", "old wording")

    resolved, _ = resolve_plan_target("presentation.replace_text", "Desktop/deck.pptx")

    assert _same(resolved, actual)
    # Discovery short-circuits on a path that exists, so re-running it is a no-op
    # and a resolved path can never be re-pointed somewhere else.
    assert _discover_readable_file(resolved) is None
    assert resolve_plan_target("presentation.replace_text", resolved) == (resolved, None)


# ── end to end: the run that failed at step 3 ────────────────────────────────
async def test_the_modifying_step_now_reaches_the_file_and_says_so(desktop: Path) -> None:
    actual = _deck(desktop / "handover" / "deck.pptx", "old wording")
    plan = build_action_plan(
        TaskRequest(text="update the wording in Desktop/deck.pptx"),
        _single_step_plan(
            "presentation.replace_text",
            "Desktop/deck.pptx",
            {"find": "old wording", "replace": "new wording"},
        ),
    )
    action = plan.actions[0]
    audit = InMemoryAuditSink()
    executor = HybridExecutor(dispatcher=build_structured_dispatcher(audit), audit=audit)

    response = await executor.execute_plan(
        plan,
        {action.action_id},
        approved_action_hashes={action.action_id: action.confirmation_hash or ""},
    )

    assert response.status == "completed"
    result = response.results[0]
    assert result.status == "succeeded"
    # The substitution stays visible, exactly as it was for the read steps.
    assert result.evidence["path_substituted"] is True
    assert _same(result.evidence["requested_path"], desktop / "deck.pptx")
    revisions = [
        event for event in audit.events if event.event_type.value == "plan_revised"
    ]
    assert len(revisions) == 1
    assert revisions[0].outcome == "path_substituted"
    # Verified by reopening the deck from disk, not from memory.
    assert result.verification is not None and result.verification.passed is True
    presentation = Presentation(str(actual))
    texts = [
        paragraph.text
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
    ]
    assert "new wording" in texts
