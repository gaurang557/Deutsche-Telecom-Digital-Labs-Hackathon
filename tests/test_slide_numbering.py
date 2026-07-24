"""Slide positions are 1-based to the planner and 0-based to the executor.

Exposing a 0-based index to the planner produced the off-by-one you would
predict: asked to update slide 3 of a three-slide deck, a live run emitted
`{"slide": 3}` and the read failed as out of range. Prompt wording had already
spelled out the 0-based rule and the model still got it wrong, so the convention
changed instead — everyone says "slide 3" for the third slide.

The conversion happens once, in `_to_executor_indexes`, at the only place a plan
becomes a `StructuredAction`. Two things are asserted here beyond the happy path:

  * the subtraction happens exactly ONCE (slide 3 returns the third slide's text,
    not the second's), and
  * `{"slide": 0}` is REFUSED rather than converted. Left alone it would become
    index -1, which Python reads as the LAST slide — a silently wrong target, and
    on a modifying action a genuinely dangerous one.

Decks are built inside `tmp_path` with a throwaway home patched over `Path.home`,
so no test touches a real Desktop.
"""

# ruff: noqa: I001

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.execution.hybrid import HybridExecutor, _to_executor_indexes
from app.planning.capabilities import find_invalid_slide_number
from app.planning.normalizer import build_action_plan
from app.schemas import DraftAction, DraftPlan, TaskRequest

#: Text placed on each slide, so the returned text identifies the slide reached.
_SLIDE_WORDING = ("first slide wording", "second slide wording", "third slide wording")


@pytest.fixture
def desktop(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    home = tmp_path / "home"
    (home / "Desktop").mkdir(parents=True)
    monkeypatch.setattr(Path, "home", classmethod(lambda cls: home))
    return home / "Desktop"


def _deck(path: Path, wording: tuple[str, ...] = _SLIDE_WORDING) -> Path:
    """A deck with one distinctly worded text box per slide."""
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    for text in wording:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        box.text_frame.paragraphs[0].text = text
    presentation.save(str(path))
    return path


def _read_plan(slide: object) -> DraftPlan:
    return DraftPlan(
        summary="I'll read that slide for you.",
        actions=[
            DraftAction(
                step_key="look",
                type="presentation.read_text",
                target="Desktop/deck.pptx",
                parameters={"slide": slide},
                description="Read the slide.",
                expected_result={"contains": "wording"},
            )
        ],
    )


async def _read(slide: object) -> object:
    plan = build_action_plan(
        TaskRequest(text="read a slide of Desktop/deck.pptx"),
        _read_plan(slide),
    )
    response = await HybridExecutor().execute_plan(plan, set())
    return response.results[0]


# ── the conversion itself ────────────────────────────────────────────────────
def test_a_slide_number_is_converted_to_an_index_exactly_once() -> None:
    assert _to_executor_indexes({"slide": 1}) == {"slide": 0}
    assert _to_executor_indexes({"slide": 3}) == {"slide": 2}


def test_other_parameters_pass_through_untouched() -> None:
    assert _to_executor_indexes({"max_chars": 3, "cell": "B2"}) == {
        "max_chars": 3,
        "cell": "B2",
    }
    # An omitted or explicitly absent slide is left exactly as it was.
    assert _to_executor_indexes({"slide": None}) == {"slide": None}


@pytest.mark.parametrize("value", [0, -1, -3])
def test_a_number_below_one_is_refused_rather_than_converted(value: int) -> None:
    with pytest.raises(ValueError, match="counts from 1"):
        _to_executor_indexes({"slide": value})


@pytest.mark.parametrize("value", [True, False, "2", 1.5, []])
def test_a_slide_that_is_not_a_whole_number_is_refused(value: object) -> None:
    with pytest.raises(ValueError, match="whole number"):
        _to_executor_indexes({"slide": value})


# ── end to end through the executor ──────────────────────────────────────────
async def test_slide_one_reaches_the_first_slide(desktop: Path) -> None:
    _deck(desktop / "deck.pptx")

    result = await _read(1)

    assert result.status == "succeeded"
    assert result.evidence["text"] == _SLIDE_WORDING[0]


async def test_the_last_slide_number_reaches_the_last_slide(desktop: Path) -> None:
    """Slide 3 of a three-slide deck: the case that failed live."""
    _deck(desktop / "deck.pptx")

    result = await _read(len(_SLIDE_WORDING))

    assert result.status == "succeeded"
    # Proves a single conversion: a double one would have returned the second.
    assert result.evidence["text"] == _SLIDE_WORDING[-1]
    assert result.evidence["text"] != _SLIDE_WORDING[-2]


async def test_one_past_the_last_slide_fails_as_out_of_range(desktop: Path) -> None:
    _deck(desktop / "deck.pptx")

    result = await _read(len(_SLIDE_WORDING) + 1)

    assert result.status == "failed"
    assert "out of range" in (result.error or "")


async def test_slide_zero_fails_and_never_resolves_to_the_last_slide(
    desktop: Path,
) -> None:
    """A -1 index would silently mean the LAST slide. It must not be reachable."""
    _deck(desktop / "deck.pptx")

    result = await _read(0)

    assert result.status == "failed"
    assert "counts from 1" in (result.error or "")
    # Nothing was read at all, least of all the final slide.
    assert result.evidence.get("text") is None


# ── the plan-time check that gives the planner a repair attempt ──────────────
def test_a_valid_slide_number_raises_no_plan_time_problem() -> None:
    assert find_invalid_slide_number({"slide": 1}) is None
    assert find_invalid_slide_number({"slide": 3}) is None
    assert find_invalid_slide_number({"max_chars": 100}) is None
    assert find_invalid_slide_number({"slide": None}) is None


@pytest.mark.parametrize("value", [0, -2])
def test_a_slide_below_one_is_reported_for_repair(value: int) -> None:
    problem = find_invalid_slide_number({"slide": value})

    assert problem is not None
    assert "counts from 1" in problem


@pytest.mark.parametrize("value", [True, "2", 2.5])
def test_a_non_integer_slide_is_reported_for_repair(value: object) -> None:
    problem = find_invalid_slide_number({"slide": value})

    assert problem is not None
    assert "whole number" in problem


def test_a_slide_bound_to_an_earlier_step_is_left_for_the_boundary() -> None:
    """Its value is unknowable at plan time, so plan time must not guess."""
    assert find_invalid_slide_number({"slide": {"$ref": "find.evidence.slide"}}) is None
    # Unchecked here, but not unchecked: the boundary still refuses a bad value.
    with pytest.raises(ValueError, match="whole number"):
        _to_executor_indexes({"slide": "resolved to nonsense"})
