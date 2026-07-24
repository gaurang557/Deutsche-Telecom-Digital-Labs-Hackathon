"""Milestone 7 — presentation.* executor + replace_text verifier.

Two layers of tests:
  1. The PresentationExecutor in isolation (each read action's success + failure
     paths, and replace_text: single-run formatting preservation, multiple
     occurrences, the `count` limit, text_not_found, save_as vs in-place).
  2. The full pipeline via the Dispatcher (AllowAllPolicy + presentation
     verifiers), proving execution AND independent re-scan verification agree
     end-to-end — including the fail-closed guard when the required verifier is
     absent.

`.pptx` fixtures are built on the fly with python-pptx in `tmp_path` (pytest
builtin), so no binary fixtures are committed and the real python-pptx
load/save path is exercised safely.
"""

import datetime as dt
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from windows_agent.contracts import Action, ActionStatus, ExecutorResult, VerificationStatus
from windows_agent.execution import ActionRegistry, Dispatcher
from windows_agent.executors.presentation_ops import (
    PresentationExecutor,
    register_presentation_executor,
)
from windows_agent.policy import AllowAllPolicy
from windows_agent.verification import (
    PresentationReplaceTextVerifier,
    VerificationRegistry,
    register_presentation_verifiers,
)


def _action(type_: str, target=None, parameters=None) -> Action:
    return Action(
        action_id="a1",
        task_id="t1",
        sequence=0,
        type=type_,
        target=str(target) if target is not None else None,
        parameters=parameters or {},
        reason="test",
    )


def _add_textbox(slide, text: str, *, bold: bool = False, italic: bool = False):
    """Add a textbox with a single-run paragraph, optionally bold/italic."""
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = box.text_frame
    tf.text = text  # one paragraph, one run
    if bold or italic:
        run = tf.paragraphs[0].runs[0]
        if bold:
            run.font.bold = True
        if italic:
            run.font.italic = True
    return box


def _make_presentation(path: Path) -> None:
    """Build a small known deck: 2 slides, textboxes (one bold run)."""
    prs = Presentation()
    blank = prs.slide_layouts[6]  # the "Blank" layout — no placeholders
    s0 = prs.slides.add_slide(blank)
    _add_textbox(s0, "Report Title")
    _add_textbox(s0, "The quick brown fox.")
    _add_textbox(s0, "Status: old value", bold=True)
    s1 = prs.slides.add_slide(blank)
    _add_textbox(s1, "The fox jumps over the fox.")
    _add_textbox(s1, "old label")
    prs.save(str(path))


def _all_text(path: Path) -> str:
    """Re-read all text-frame text from a deck (for assertions)."""
    prs = Presentation(str(path))
    pieces = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                pieces.append(shape.text_frame.text)
    return "\n".join(pieces)


# ── PresentationExecutor unit tests: read-only ─────────────────────────────
async def test_slide_count(tmp_path):
    ppt = tmp_path / "deck.pptx"
    _make_presentation(ppt)
    ex = PresentationExecutor()
    res = await ex.execute(_action("presentation.slide_count", ppt))
    assert res.success is True
    assert res.evidence["slide_count"] == 2


async def test_read_text(tmp_path):
    ppt = tmp_path / "deck.pptx"
    _make_presentation(ppt)
    ex = PresentationExecutor()
    res = await ex.execute(_action("presentation.read_text", ppt))
    assert res.success is True
    text = res.evidence["text"]
    assert "Report Title" in text
    assert "The quick brown fox." in text
    assert "old label" in text
    # Empty paragraphs are not joined, so no doubled blank lines.
    assert "\n\n" not in text
    assert res.evidence["truncated"] is False
    assert res.evidence["slides_read"] == 2


async def test_read_text_truncation(tmp_path):
    ppt = tmp_path / "deck.pptx"
    _make_presentation(ppt)
    ex = PresentationExecutor()
    res = await ex.execute(_action("presentation.read_text", ppt, {"max_chars": 5}))
    assert res.success is True
    assert res.evidence["truncated"] is True
    assert len(res.evidence["text"]) == 5


async def test_read_text_single_slide(tmp_path):
    ppt = tmp_path / "deck.pptx"
    _make_presentation(ppt)
    ex = PresentationExecutor()
    res = await ex.execute(_action("presentation.read_text", ppt, {"slide": 1}))
    assert res.success is True
    text = res.evidence["text"]
    # Only slide index 1 was read.
    assert res.evidence["slides_read"] == 1
    assert "fox jumps over" in text
    assert "old label" in text
    assert "Report Title" not in text
    assert "quick brown" not in text


async def test_read_text_slide_out_of_range(tmp_path):
    ppt = tmp_path / "deck.pptx"
    _make_presentation(ppt)
    ex = PresentationExecutor()
    res = await ex.execute(_action("presentation.read_text", ppt, {"slide": 99}))
    assert res.success is False
    assert res.error.code == "slide_out_of_range"


async def test_get_metadata(tmp_path):
    ppt = tmp_path / "deck.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.core_properties.title = "My Deck"
    prs.core_properties.author = "Ada"
    prs.core_properties.subject = "Testing"
    prs.core_properties.keywords = "alpha,beta"
    prs.core_properties.created = dt.datetime(2024, 1, 2, 3, 4, 5)
    prs.core_properties.last_modified_by = "Grace"
    prs.save(str(ppt))

    ex = PresentationExecutor()
    res = await ex.execute(_action("presentation.get_metadata", ppt))
    assert res.success is True
    meta = res.evidence["metadata"]
    assert meta["title"] == "My Deck"
    assert meta["author"] == "Ada"
    assert meta["subject"] == "Testing"
    assert meta["keywords"] == "alpha,beta"
    assert meta["last_modified_by"] == "Grace"
    # Datetime is an ISO string; compare on the stable prefix.
    assert isinstance(meta["created"], str)
    assert meta["created"].startswith("2024-01-02T03:04:05")
    # Only the seven documented fields are surfaced.
    assert set(meta) == {
        "title",
        "author",
        "subject",
        "keywords",
        "created",
        "modified",
        "last_modified_by",
    }


async def test_find_per_slide_counts(tmp_path):
    ppt = tmp_path / "deck.pptx"
    _make_presentation(ppt)
    ex = PresentationExecutor()
    res = await ex.execute(_action("presentation.find", ppt, {"query": "fox"}))
    assert res.success is True
    # slide 0 "The quick brown fox." (1) and slide 1 "The fox ... the fox." (2).
    counts = {m["count"] for m in res.evidence["matches"]}
    slides = {m["slide_index"] for m in res.evidence["matches"]}
    assert res.evidence["total_matches"] == 3
    assert counts == {1, 2}
    assert slides == {0, 1}
    assert res.evidence["truncated"] is False


async def test_find_empty_query_fails(tmp_path):
    ppt = tmp_path / "deck.pptx"
    _make_presentation(ppt)
    ex = PresentationExecutor()
    res = await ex.execute(_action("presentation.find", ppt, {"query": ""}))
    assert res.success is False
    assert res.error.code == "invalid_parameters"


# ── PresentationExecutor unit tests: replace_text ──────────────────────────
async def test_replace_text_preserves_run_formatting(tmp_path):
    """A single-run replacement keeps that run's bold/italic formatting."""
    ppt = tmp_path / "fmt.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, "Status: old value", bold=True, italic=True)
    prs.save(str(ppt))

    ex = PresentationExecutor()
    res = await ex.execute(
        _action("presentation.replace_text", ppt, {"find": "old", "replace": "new"})
    )
    assert res.success is True
    assert res.evidence["replacements"] == 1
    assert res.evidence["save_as"] is False
    assert res.evidence["output_path"] == str(ppt)

    reopened = Presentation(str(ppt))
    tf = reopened.slides[0].shapes[0].text_frame
    para = tf.paragraphs[0]
    assert para.text == "Status: new value"
    # Formatting survives on the replaced run.
    target = next(r for r in para.runs if "new" in r.text)
    assert target.font.bold is True
    assert target.font.italic is True


async def test_replace_text_multiple_occurrences(tmp_path):
    ppt = tmp_path / "deck.pptx"
    _make_presentation(ppt)
    ex = PresentationExecutor()
    # "old" appears on slide 0 ("Status: old value") and slide 1 ("old label").
    res = await ex.execute(
        _action("presentation.replace_text", ppt, {"find": "old", "replace": "NEW"})
    )
    assert res.success is True
    assert res.evidence["replacements"] == 2
    all_text = _all_text(ppt)
    assert "old" not in all_text
    assert all_text.count("NEW") == 2


async def test_replace_text_count_limit(tmp_path):
    ppt = tmp_path / "count.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, "foo foo foo")
    prs.save(str(ppt))
    ex = PresentationExecutor()
    res = await ex.execute(
        _action("presentation.replace_text", ppt, {"find": "foo", "replace": "bar", "count": 2})
    )
    assert res.success is True
    assert res.evidence["replacements"] == 2
    reopened = Presentation(str(ppt))
    assert reopened.slides[0].shapes[0].text_frame.paragraphs[0].text == "bar bar foo"


async def test_replace_text_not_found(tmp_path):
    ppt = tmp_path / "deck.pptx"
    _make_presentation(ppt)
    ex = PresentationExecutor()
    res = await ex.execute(
        _action("presentation.replace_text", ppt, {"find": "absent-string", "replace": "x"})
    )
    assert res.success is False
    assert res.error.code == "text_not_found"


async def test_replace_text_empty_find_fails(tmp_path):
    ppt = tmp_path / "deck.pptx"
    _make_presentation(ppt)
    ex = PresentationExecutor()
    res = await ex.execute(
        _action("presentation.replace_text", ppt, {"find": "", "replace": "x"})
    )
    assert res.success is False
    assert res.error.code == "invalid_parameters"


async def test_replace_text_save_as_leaves_original_untouched(tmp_path):
    src = tmp_path / "src.pptx"
    out = tmp_path / "out.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, "keep old text")
    prs.save(str(src))

    ex = PresentationExecutor()
    res = await ex.execute(
        _action("presentation.replace_text", src, {"find": "old", "replace": "new", "save_as": str(out)})
    )
    assert res.success is True
    assert res.evidence["save_as"] is True
    assert res.evidence["output_path"] == str(out)
    # Original untouched; new file has the correction.
    assert "keep old text" in _all_text(src)
    assert "keep new text" in _all_text(out)


async def test_replace_text_in_place_edits_original(tmp_path):
    ppt = tmp_path / "inplace.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, "edit old here")
    prs.save(str(ppt))
    ex = PresentationExecutor()
    res = await ex.execute(
        _action("presentation.replace_text", ppt, {"find": "old", "replace": "new"})
    )
    assert res.success is True
    assert "edit new here" in _all_text(ppt)


async def test_replace_text_save_as_existing_fails_without_overwrite(tmp_path):
    src = tmp_path / "src.pptx"
    out = tmp_path / "existing.pptx"
    Presentation().save(str(out))  # pre-existing target
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, "has old text")
    prs.save(str(src))

    ex = PresentationExecutor()
    res = await ex.execute(
        _action("presentation.replace_text", src, {"find": "old", "replace": "new", "save_as": str(out)})
    )
    assert res.success is False
    assert res.error.code == "output_exists"


async def test_replace_text_cross_run_fallback_collapses_formatting(tmp_path):
    """A match spanning runs collapses to the first run (documented limitation)."""
    ppt = tmp_path / "crossrun.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    para = box.text_frame.paragraphs[0]
    r1 = para.add_run()
    r1.text = "Hel"
    r1.font.bold = True
    r2 = para.add_run()
    r2.text = "lo world"  # not bold
    prs.save(str(ppt))

    ex = PresentationExecutor()
    res = await ex.execute(
        _action("presentation.replace_text", ppt, {"find": "Hello", "replace": "Goodbye"})
    )
    assert res.success is True
    assert res.evidence["replacements"] == 1
    reopened = Presentation(str(ppt))
    assert reopened.slides[0].shapes[0].text_frame.paragraphs[0].text == "Goodbye world"


# ── Error paths ─────────────────────────────────────────────────────────────
async def test_read_missing_file(tmp_path):
    ex = PresentationExecutor()
    res = await ex.execute(_action("presentation.read_text", tmp_path / "nope.pptx"))
    assert res.success is False
    assert res.error.code == "file_not_found"


async def test_non_pptx_file_fails(tmp_path):
    fake = tmp_path / "data.txt"
    fake.write_text("not a presentation")
    ex = PresentationExecutor()
    res = await ex.execute(_action("presentation.read_text", fake))
    assert res.success is False
    assert res.error.code == "not_a_presentation"


# ── Verifier unit tests ────────────────────────────────────────────────────
async def test_replace_verifier_passes(tmp_path):
    ppt = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, "correct the old spelling")
    prs.save(str(ppt))
    ex = PresentationExecutor()
    action = _action("presentation.replace_text", ppt, {"find": "old", "replace": "new"})
    res = await ex.execute(action)
    assert res.success is True
    vr = await PresentationReplaceTextVerifier().verify(action, res)
    assert vr.status == VerificationStatus.PASSED
    assert vr.observed["replace_count"] >= 1
    assert vr.observed["find_count"] == 0


async def test_replace_verifier_fails_when_replacement_absent(tmp_path):
    """Evidence claims a replacement that is not actually present → FAILED."""
    ppt = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, "nothing to see")
    prs.save(str(ppt))
    action = _action("presentation.replace_text", ppt, {"find": "old", "replace": "new"})
    # Fabricated evidence: says it replaced once, but "new" is nowhere on disk.
    fake = ExecutorResult(
        success=True,
        evidence={"output_path": str(ppt), "find": "old", "replace": "new", "replacements": 1},
    )
    vr = await PresentationReplaceTextVerifier().verify(action, fake)
    assert vr.status == VerificationStatus.FAILED


async def test_replace_verifier_fails_when_reverted(tmp_path):
    ppt = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, "fix the old bug")
    prs.save(str(ppt))
    ex = PresentationExecutor()
    action = _action("presentation.replace_text", ppt, {"find": "old", "replace": "new"})
    res = await ex.execute(action)
    # Externally revert the deck so re-observation disagrees with intent.
    reverted = Presentation()
    rslide = reverted.slides.add_slide(reverted.slide_layouts[6])
    _add_textbox(rslide, "fix the old bug")
    reverted.save(str(ppt))
    vr = await PresentationReplaceTextVerifier().verify(action, res)
    assert vr.status == VerificationStatus.FAILED


# ── End-to-end via the Dispatcher (policy + execution + verification) ───────
def _pipeline():
    reg = ActionRegistry()
    register_presentation_executor(reg)
    vreg = VerificationRegistry()
    register_presentation_verifiers(vreg)
    return Dispatcher(reg, AllowAllPolicy(), verification=vreg)


async def test_pipeline_replace_text_success_and_verified(tmp_path):
    ppt = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, "please correct old text here")
    prs.save(str(ppt))
    disp = _pipeline()
    result = await disp.dispatch(
        _action("presentation.replace_text", ppt, {"find": "old", "replace": "new"})
    )
    assert result.status == ActionStatus.SUCCESS
    assert result.verification.status == VerificationStatus.PASSED
    assert "please correct new text here" in _all_text(ppt)


async def test_pipeline_read_text_skips_verification(tmp_path):
    ppt = tmp_path / "deck.pptx"
    _make_presentation(ppt)
    disp = _pipeline()
    result = await disp.dispatch(_action("presentation.read_text", ppt))
    assert result.status == ActionStatus.SUCCESS
    # Read-only → no verifier registered → SKIPPED.
    assert result.verification.status == VerificationStatus.SKIPPED


async def test_pipeline_replace_text_fails_closed_without_verifier(tmp_path):
    """A required-verification action with NO registered verifier fails closed."""
    ppt = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, "leave old alone")
    prs.save(str(ppt))
    reg = ActionRegistry()
    register_presentation_executor(reg)
    # Deliberately DO NOT register presentation verifiers.
    disp = Dispatcher(reg, AllowAllPolicy(), verification=VerificationRegistry())
    result = await disp.dispatch(
        _action("presentation.replace_text", ppt, {"find": "old", "replace": "new"})
    )
    assert result.status == ActionStatus.FAILED
    assert result.error.code == "verifier_missing"
    # The executor never ran, so the deck is untouched.
    assert "leave old alone" in _all_text(ppt)


async def test_register_presentation_executor_marks_verification_requirement():
    reg = ActionRegistry()
    register_presentation_executor(reg)
    # The one modifying action requires verification; reads do not (the
    # fail-closed contract keys off this deterministic metadata).
    assert reg.requires_verification("presentation.replace_text") is True
    for read_type in (
        "presentation.slide_count",
        "presentation.get_metadata",
        "presentation.read_text",
        "presentation.find",
    ):
        assert reg.requires_verification(read_type) is False


async def test_register_presentation_executor_covers_all_types():
    reg = ActionRegistry()
    register_presentation_executor(reg)
    for t in (
        "presentation.slide_count",
        "presentation.get_metadata",
        "presentation.read_text",
        "presentation.find",
        "presentation.replace_text",
    ):
        assert reg.has_action(t)
