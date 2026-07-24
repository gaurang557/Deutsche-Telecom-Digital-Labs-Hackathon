"""`PresentationExecutor` — the `.pptx` PowerPoint executor (Milestone 7).

WHAT THIS IS
------------
A single executor that performs the `presentation.*` semantic actions against
`.pptx` presentations via **python-pptx** (`from pptx import Presentation`) — a
structured presentation API, far more reliable than scraping a viewer's GUI (see
ARCHITECTURE §3 executor-preference order). Like `FileExecutor` / `PdfExecutor`
/ `SpreadsheetExecutor` / `DocumentExecutor`, it is deliberately "dumb" about
safety: it does NOT decide permissions or ask for confirmation. That is the
Policy Engine's job (the dispatcher only runs this executor after an ALLOW
decision). See `executors/base.py`.

ACTION VOCABULARY (this milestone)
----------------------------------
Read-only (no side effects, verifier SKIPPED):
  * presentation.slide_count  — how many slides does the deck have?
  * presentation.get_metadata — the deck's core properties (title/author/…).
  * presentation.read_text    — text of one slide or all slides (bounded).
  * presentation.find         — per-slide case-sensitive substring hit counts.

Modifying (verified by re-observation in
`verification/presentation_verifiers.py`):
  * presentation.replace_text — replace occurrences of `find` with `replace`,
                                **preserving formatting**, across every slide's
                                text-frame runs. The headline M7 capability:
                                correcting a deck without disturbing formatting.

PARAMETER CONVENTIONS
---------------------
`action.target` is the PRIMARY path (the presentation `.pptx`). Everything else
lives in `action.parameters`, e.g. {"max_chars": 5000}, {"slide": 2},
{"query": "foo"}, {"find": "old", "replace": "new", "count": 2,
"save_as": "out.pptx"}. `slide` is a **0-based** slide index selecting a single
slide (validated against the real slide count; fail closed on range). Only
`.pptx` is supported (the legacy binary `.ppt` format is not).

FORMATTING PRESERVATION — approach + its documented limitation
--------------------------------------------------------------
A `.pptx` paragraph (inside a shape's text frame) is a sequence of *runs*, and
each run carries its own formatting (bold/italic/font/…). PowerPoint freely
splits a logical piece of text across several runs, so a search string may live
entirely inside one run or straddle a run boundary.

We handle the two cases pragmatically (identically to `document_ops`):
  * **Match within a single run (the clean case):** we replace the substring
    *inside that run's text only*, so the run — and therefore all of its
    formatting — is preserved exactly. This is the common case and the one we
    assert on in tests.
  * **Match spanning multiple runs (the fallback):** there is no single run
    whose formatting is "the" formatting of the match, so we rebuild the
    paragraph text across its runs, apply the replacement, write the whole
    result into the paragraph's FIRST run, and clear the remaining runs.
    **LIMITATION:** a cross-run replacement therefore collapses the affected
    text to the first run's formatting. This is an accepted, documented
    trade-off (a fully faithful cross-run edit would require splitting/merging
    runs at XML level, which is out of scope for M7).

WHY ASYNC + to_thread
---------------------
The executor contract is async, but python-pptx's load/save calls are blocking
(they parse/serialise a zip of XML on disk). Each operation is dispatched to a
worker thread via `asyncio.to_thread` instead of blocking the event loop —
exactly as file_ops/pdf_ops/spreadsheet_ops/document_ops do for their blocking
I/O. python-pptx reads the whole package into memory on open and keeps no file
handle, so a loaded presentation is simply discarded (no explicit close).

SAFETY / BOUNDING
-----------------
  * `presentation.read_text` is capped at `_DEFAULT_TEXT_CHAR_CAP` characters
    (override with `max_chars`) and `presentation.find` at
    `_DEFAULT_SEARCH_RESULTS` matching-slide entries, so a huge deck can never
    be slurped into memory / evidence; `truncated=true` marks a clipped result.
    The dispatcher additionally bounds evidence.
  * `presentation.replace_text` with an empty `find` is rejected; a `find` that
    is absent everywhere fails with `text_not_found` (0 replacements is reported
    as an ERROR, so the planner learns the correction did not apply and nothing
    is written).
  * `save_as` writes the result to a NEW path (the original is untouched);
    without it the executor edits the deck IN PLACE (overwriting the original).
    A `save_as` that would clobber an existing DIFFERENT file fails closed with
    `output_exists` unless `overwrite=true`.
  * Expected errors are returned as `ExecutorResult(success=False, error=...)` —
    the executor never raises for ordinary failures. (Unexpected exceptions are
    still contained by the dispatcher.)

RISK (documentation only — the executor NEVER sets risk)
--------------------------------------------------------
  * reads (`slide_count`/`get_metadata`/`read_text`/`find`) → `RiskLevel.NONE`.
  * `presentation.replace_text` editing IN PLACE overwrites the original → `HIGH`.
  * `presentation.replace_text` with `save_as` writes a new file → `MEDIUM`.
Risk is assigned by the deterministic policy (M12); it is documented here but
never decided by this executor.
"""

from __future__ import annotations

import asyncio
import datetime as _dt
from pathlib import Path
from typing import Any, Iterator

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..contracts import Action, ActionError, ErrorCode, ExecutorResult
from .base import BaseExecutor

# Domain-specific error codes. `ActionError.code` is a free string (see
# contracts/error.py); we use stable presentation-specific codes here and fall
# back to the shared ErrorCode values where they fit.
ERR_FILE_NOT_FOUND = "file_not_found"
ERR_NOT_A_PRESENTATION = "not_a_presentation"
ERR_TEXT_NOT_FOUND = "text_not_found"
ERR_OUTPUT_EXISTS = "output_exists"
ERR_INVALID_PARAMS = "invalid_parameters"
ERR_SLIDE_OUT_OF_RANGE = "slide_out_of_range"

#: Default cap on extracted text length so a huge deck is never read into
#: evidence wholesale. Overridable per call via `max_chars`.
_DEFAULT_TEXT_CHAR_CAP = 20_000
#: Default/limit on the number of per-slide match entries `presentation.find`
#: returns, so a query hitting every slide cannot bloat evidence.
_DEFAULT_SEARCH_RESULTS = 100
#: Sentinel used when `count` is not supplied (replace ALL occurrences). A deck
#: will never contain this many occurrences.
_UNLIMITED = 1_000_000_000

#: Every action type this executor handles and its deterministic verification
#: requirement.
PRESENTATION_ACTION_REQUIREMENTS: dict[str, bool] = {
    "presentation.slide_count": False,
    "presentation.get_metadata": False,
    "presentation.read_text": False,
    "presentation.find": False,
    "presentation.replace_text": True,
}
PRESENTATION_ACTION_TYPES: tuple[str, ...] = tuple(PRESENTATION_ACTION_REQUIREMENTS)

#: core_properties attributes we surface. Datetimes become ISO strings; empty
#: strings become null; everything else passes through.
_METADATA_FIELDS = (
    "title",
    "author",
    "subject",
    "keywords",
    "created",
    "modified",
    "last_modified_by",
)


def _err(code: str, message: str, *, retryable: bool = False) -> ExecutorResult:
    return ExecutorResult(
        success=False,
        error=ActionError(code=code, message=message, retryable=retryable),
    )


def _normalize_meta(value: Any) -> Any:
    """Coerce a core-property value to a JSON-serialisable primitive.

    Datetimes → ISO-8601 strings, empty strings → None, everything else
    (non-empty strings, ints) passes through unchanged.
    """
    if value is None:
        return None
    if isinstance(value, (_dt.datetime, _dt.date, _dt.time)):
        return value.isoformat()
    if isinstance(value, str):
        return value or None
    return value


def _iter_shapes(shapes: Any) -> Iterator[Any]:
    """Yield every shape, recursing into group shapes.

    PowerPoint lets authors nest shapes inside groups; a group itself has no
    text frame but its children do, so we descend into groups to reach all the
    text a plain top-level walk would miss.
    """
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _slide_text_frames(slide: Any) -> Iterator[Any]:
    """Yield the text frame of every text-bearing shape on one slide."""
    for shape in _iter_shapes(slide.shapes):
        if shape.has_text_frame:
            yield shape.text_frame


def _slide_paragraphs(slide: Any) -> Iterator[Any]:
    """Yield every paragraph on one slide (across all its text frames)."""
    for text_frame in _slide_text_frames(slide):
        yield from text_frame.paragraphs


def _all_paragraphs(prs: Any) -> Iterator[Any]:
    """Yield every paragraph we edit/scan across the whole deck."""
    for slide in prs.slides:
        yield from _slide_paragraphs(slide)


def _slide_text(slide: Any) -> str:
    """Join the non-empty paragraph texts of one slide with newlines."""
    return "\n".join(p.text for p in _slide_paragraphs(slide) if p.text)


def _replace_in_paragraph(paragraph: Any, find: str, replace: str, remaining: int) -> int:
    """Replace up to `remaining` occurrences of `find` in one paragraph.

    Returns the number of replacements performed. See the module docstring for
    the run-level (formatting-preserving) vs cross-run (fallback) strategy.
    """
    if remaining <= 0 or not paragraph.runs:
        return 0

    original_full = "".join(run.text for run in paragraph.runs)
    full_count = original_full.count(find)
    if full_count == 0:
        return 0
    # Occurrences wholly contained in a single run (non-overlapping counting,
    # so this is always <= full_count).
    within_run_count = sum(run.text.count(find) for run in paragraph.runs)

    # Clean case — every match lives inside a single run, so we can replace in
    # place and each run keeps ALL of its formatting. `str.replace` operates on
    # the original run text and never re-scans its own output, so this stays
    # correct even when `replace` contains `find` (e.g. "old" → "older").
    if within_run_count == full_count:
        done = 0
        for run in paragraph.runs:
            if remaining - done <= 0:
                break
            occurrences = run.text.count(find)
            if occurrences:
                allowed = min(occurrences, remaining - done)
                run.text = run.text.replace(find, replace, allowed)
                done += allowed
        return done

    # Fallback — at least one match straddles a run boundary. Replace against the
    # ORIGINAL joined text (so replacement output is never re-matched) and write
    # the whole result into the first run, clearing the rest. This collapses the
    # paragraph to the first run's formatting — the documented M7 limitation, and
    # only for paragraphs that actually contain a cross-run match.
    allowed = min(full_count, remaining)
    paragraph.runs[0].text = original_full.replace(find, replace, allowed)
    for run in paragraph.runs[1:]:
        run.text = ""
    return allowed


class PresentationExecutor(BaseExecutor):
    """Executes the `presentation.*` action vocabulary via python-pptx (.pptx)."""

    name = "presentation"

    async def execute(self, action: Action) -> ExecutorResult:
        # Route on the semantic type. The registry only sends us presentation.*
        # types it was told to, but we still guard against an unexpected one.
        handler = {
            "presentation.slide_count": self._slide_count,
            "presentation.get_metadata": self._get_metadata,
            "presentation.read_text": self._read_text,
            "presentation.find": self._find,
            "presentation.replace_text": self._replace_text,
        }.get(action.type)
        if handler is None:
            return _err(
                ErrorCode.NOT_IMPLEMENTED.value,
                f"PresentationExecutor does not handle action type {action.type!r}",
            )
        # Blocking python-pptx load/save runs off the event loop.
        return await asyncio.to_thread(handler, action)

    # ── helpers ────────────────────────────────────────────────────────────
    @staticmethod
    def _require_target(action: Action) -> Path | None:
        return Path(action.target) if action.target else None

    @staticmethod
    def _param(action: Action, key: str, default: Any = None) -> Any:
        return action.parameters.get(key, default)

    @staticmethod
    def _as_index(value: Any) -> int | None:
        """Coerce a slide parameter to a non-negative int, or None if invalid."""
        if isinstance(value, bool):  # bool is an int subclass; reject it explicitly
            return None
        if not isinstance(value, int):
            return None
        return value if value >= 0 else None

    def _open_presentation(self, path: Path | None) -> tuple[Any, ExecutorResult | None]:
        """Open + guard an existing `.pptx`, returning (prs, None) or
        (None, error_result).

        Fails closed on a missing path, a non-file target, a non-`.pptx`
        extension, and anything python-pptx cannot parse as a presentation.
        python-pptx keeps no open file handle, so the caller just discards the
        returned presentation.
        """
        if path is None:
            return None, _err(ERR_INVALID_PARAMS, "presentation action requires a target path")
        if not path.exists():
            return None, _err(ERR_FILE_NOT_FOUND, f"File not found: {path}")
        if not path.is_file():
            return None, _err(ERR_NOT_A_PRESENTATION, f"Not a file: {path}")
        if path.suffix.lower() != ".pptx":
            return None, _err(ERR_NOT_A_PRESENTATION, f"Not a .pptx presentation: {path}")
        try:
            prs = Presentation(str(path))
        except Exception as exc:  # python-pptx raises various zip/XML/format errors
            return None, _err(ERR_NOT_A_PRESENTATION, f"Cannot open as .pptx: {path}: {exc}")
        return prs, None

    # ── read-only operations ────────────────────────────────────────────────
    def _slide_count(self, action: Action) -> ExecutorResult:
        path = self._require_target(action)
        prs, err = self._open_presentation(path)
        if err is not None:
            return err
        return ExecutorResult(
            success=True,
            evidence={"path": str(path), "slide_count": len(prs.slides)},
        )

    def _get_metadata(self, action: Action) -> ExecutorResult:
        path = self._require_target(action)
        prs, err = self._open_presentation(path)
        if err is not None:
            return err
        props = prs.core_properties
        metadata = {field: _normalize_meta(getattr(props, field, None)) for field in _METADATA_FIELDS}
        return ExecutorResult(
            success=True,
            evidence={"path": str(path), "metadata": metadata},
        )

    def _read_text(self, action: Action) -> ExecutorResult:
        path = self._require_target(action)
        prs, err = self._open_presentation(path)
        if err is not None:
            return err

        max_chars = self._param(action, "max_chars", _DEFAULT_TEXT_CHAR_CAP)
        try:
            max_chars = int(max_chars)
        except (TypeError, ValueError):
            return _err(ERR_INVALID_PARAMS, f"max_chars must be an integer, got {max_chars!r}")
        if max_chars <= 0:
            max_chars = _DEFAULT_TEXT_CHAR_CAP

        slides = prs.slides
        slide_count = len(slides)
        # A `slide` parameter selects exactly one 0-based slide; else read all.
        slide_param = self._param(action, "slide")
        if slide_param is not None:
            index = self._as_index(slide_param)
            if index is None:
                return _err(ERR_INVALID_PARAMS, f"slide must be a non-negative integer, got {slide_param!r}")
            if index >= slide_count:
                return _err(
                    ERR_SLIDE_OUT_OF_RANGE,
                    f"slide {index} out of range (presentation has {slide_count} slide(s))",
                )
            selected = [slides[index]]
        else:
            selected = list(slides)

        # Join only the non-empty paragraphs across the selected slides, capping
        # the total length. `slides_read` counts the slides we actually visited
        # (including the one on which truncation occurred).
        parts: list[str] = []
        total = 0
        truncated = False
        slides_read = 0
        for slide in selected:
            slides_read += 1
            stop = False
            for para in _slide_paragraphs(slide):
                text = para.text
                if not text:
                    continue
                piece = ("\n" if parts else "") + text
                if total + len(piece) > max_chars:
                    parts.append(piece[: max_chars - total])
                    truncated = True
                    stop = True
                    break
                parts.append(piece)
                total += len(piece)
            if stop:
                break

        return ExecutorResult(
            success=True,
            evidence={
                "path": str(path),
                "text": "".join(parts),
                "slides_read": slides_read,
                "truncated": truncated,
            },
        )

    def _find(self, action: Action) -> ExecutorResult:
        path = self._require_target(action)
        query = self._param(action, "query")
        if not isinstance(query, str) or query == "":
            # Open validity is irrelevant if the query is unusable; fail fast,
            # but still surface a missing/broken file first for a clearer error.
            prs, err = self._open_presentation(path)
            if err is not None:
                return err
            return _err(ERR_INVALID_PARAMS, "presentation.find requires a non-empty parameters.query")

        prs, err = self._open_presentation(path)
        if err is not None:
            return err

        max_results = self._param(action, "max_results", _DEFAULT_SEARCH_RESULTS)
        try:
            max_results = int(max_results)
        except (TypeError, ValueError):
            max_results = _DEFAULT_SEARCH_RESULTS
        if max_results <= 0:
            max_results = _DEFAULT_SEARCH_RESULTS

        matches: list[dict[str, int]] = []
        total_matches = 0
        truncated = False
        for index, slide in enumerate(prs.slides):
            count = _slide_text(slide).count(query)
            if not count:
                continue
            if len(matches) >= max_results:
                # More slides match than we will report — mark bounded.
                truncated = True
                break
            matches.append({"slide_index": index, "count": count})
            total_matches += count

        return ExecutorResult(
            success=True,
            evidence={
                "path": str(path),
                "matches": matches,
                "total_matches": total_matches,
                "truncated": truncated,
            },
        )

    # ── modifying operation ─────────────────────────────────────────────────
    def _replace_text(self, action: Action) -> ExecutorResult:
        path = self._require_target(action)
        find = self._param(action, "find")
        replace = self._param(action, "replace", "")

        if not isinstance(find, str) or find == "":
            return _err(ERR_INVALID_PARAMS, "presentation.replace_text requires a non-empty parameters.find")
        if not isinstance(replace, str):
            return _err(ERR_INVALID_PARAMS, "presentation.replace_text parameters.replace must be a string")

        # Optional count limit: when supplied it must be a positive int.
        count_raw = self._param(action, "count")
        if count_raw is None:
            remaining = _UNLIMITED
        else:
            if isinstance(count_raw, bool) or not isinstance(count_raw, int) or count_raw <= 0:
                return _err(ERR_INVALID_PARAMS, f"count must be a positive integer, got {count_raw!r}")
            remaining = count_raw

        # Resolve the output path: save_as → a NEW file; otherwise edit in place.
        save_as = self._param(action, "save_as")
        overwrite = bool(self._param(action, "overwrite", False))
        if save_as is not None:
            if not isinstance(save_as, str) or not save_as.strip():
                return _err(ERR_INVALID_PARAMS, "save_as must be a non-empty path string")
            output_path = Path(save_as)
            if output_path.suffix.lower() != ".pptx":
                return _err(ERR_INVALID_PARAMS, f"save_as must be a .pptx path: {output_path}")
            if not output_path.parent.exists():
                return _err(ERR_FILE_NOT_FOUND, f"Directory does not exist: {output_path.parent}")
        else:
            output_path = path  # in-place edit (overwrites the original)

        prs, err = self._open_presentation(path)
        if err is not None:
            return err

        # Refuse to clobber a DIFFERENT existing file via save_as unless allowed
        # (mirrors document.replace_text / file.write_text's fail-closed guard).
        if (
            save_as is not None
            and output_path.exists()
            and output_path.resolve() != path.resolve()
            and not overwrite
        ):
            return _err(
                ERR_OUTPUT_EXISTS,
                f"save_as target already exists (set overwrite=true to replace): {output_path}",
            )

        total = 0
        for para in _all_paragraphs(prs):
            if remaining - total <= 0:
                break
            total += _replace_in_paragraph(para, find, replace, remaining - total)

        # 0 replacements is an ERROR (text_not_found), so the planner learns the
        # correction did not apply and nothing is written.
        if total == 0:
            return _err(ERR_TEXT_NOT_FOUND, f"Text not found in presentation: {find!r}")

        prs.save(str(output_path))
        return ExecutorResult(
            success=True,
            evidence={
                "path": str(path),
                "output_path": str(output_path),
                "find": find,
                "replace": replace,
                "replacements": total,
                "save_as": save_as is not None,
            },
            side_effects=[{"type": "presentation.text_replaced", "target": str(output_path)}],
        )


def register_presentation_executor(
    registry, executor: PresentationExecutor | None = None, *, override: bool = False
) -> PresentationExecutor:
    """Register a single `PresentationExecutor` for every `presentation.*` type.

    Returns the executor instance so callers can reuse it. One instance handles
    all presentation operations (it is stateless), matching the "one executor,
    many action types" pattern noted in `executors/base.py`.
    """
    executor = executor or PresentationExecutor()
    for action_type, requires_verification in PRESENTATION_ACTION_REQUIREMENTS.items():
        registry.register_action(
            action_type,
            executor,
            requires_verification=requires_verification,
            override=override,
        )
    return executor
