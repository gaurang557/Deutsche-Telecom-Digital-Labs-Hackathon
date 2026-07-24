r"""Interactive manual tester for the `presentation.*` functionality (Milestone 7).

This is a HUMAN-DRIVEN harness (not an automated test). It wires up the real
execution pipeline exactly as the agent will use it —

    ActionRegistry (+ PresentationExecutor)
      → Dispatcher (+ AllowAllPolicy, VerificationRegistry (+ presentation
        verifiers), InMemoryAuditSink)

— then lets you type presentation commands and shows, for each one:
  * the resulting ActionStatus,
  * the bounded evidence the executor returned (slide count, text, metadata,
    matches, replacement counts),
  * the verification result (SKIPPED for reads; an INDEPENDENT re-scan PASS/FAIL
    for `replace`/`replaceas`), and
  * the ordered audit trail emitted around the action.

Relative paths you type are resolved inside a throwaway `sandbox/` directory;
absolute paths are used as-is. Only `.pptx` presentations are supported.

To get started with no deck of your own, run `sample` to generate a small
`sample.pptx` (a couple of slides with text including bold text, plus core
properties) inside the sandbox, then exercise the commands against it.

RUN (from the windows-agent/ folder, using the project venv):

    & "..\.venv\Scripts\python.exe" tools\manual_presentation_test.py
    # or point at a custom sandbox:
    & "..\.venv\Scripts\python.exe" tools\manual_presentation_test.py --sandbox C:\tmp\wa

Type `help` at the prompt for the command list, `quit` to exit.

NOTE ON POLICY: this harness uses AllowAllPolicy, so every action is authorized
(the real deterministic policy arrives in M12). The point here is to exercise
PresentationExecutor behaviour + the re-scan replace verifier, which is what M7
delivers. `replace` edits the file IN PLACE (HIGH risk); `replaceas` writes a
NEW file and leaves the original untouched (MEDIUM risk).
"""

from __future__ import annotations

import argparse
import asyncio
import shlex
import sys
import uuid
from pathlib import Path

# Make the module importable when run as a plain script (sys.path[0] is tools/).
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation  # also used to synthesise the sample deck  # noqa: E402
from pptx.util import Inches  # noqa: E402

from windows_agent.audit import InMemoryAuditSink  # noqa: E402
from windows_agent.contracts import Action  # noqa: E402
from windows_agent.execution import ActionRegistry, Dispatcher  # noqa: E402
from windows_agent.executors import register_presentation_executor  # noqa: E402
from windows_agent.policy import AllowAllPolicy  # noqa: E402
from windows_agent.verification import VerificationRegistry, register_presentation_verifiers  # noqa: E402


HELP = """
Commands (paths are relative to the sandbox unless absolute; .pptx only):
  sample [name]                    generate a sample deck (default: sample.pptx)
  ls                               list decks in the sandbox
  slides <path>                    presentation.slide_count — how many slides
  read <path> [max_chars]          presentation.read_text   — the deck's text
  meta <path>                      presentation.get_metadata — core properties
  find <path> <query>              presentation.find        — per-slide hit counts
  replace <path> <find> <replace> [count]
                                   presentation.replace_text — edit IN PLACE (HIGH)
  replaceas <path> <find> <replace> <out> [count]
                                   presentation.replace_text — write NEW file (MEDIUM)
  help                             show this help
  quit / exit                      leave

Try this quick scenario:
  sample
  read sample.pptx
  find sample.pptx old
  replace sample.pptx old new
  read sample.pptx
"""

_SAMPLE_HEADING = "Quarterly Report"


def _add_textbox(slide, text: str, *, bold: bool = False):
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    box.text_frame.text = text
    if bold:
        box.text_frame.paragraphs[0].runs[0].font.bold = True
    return box


def _make_sample(path: Path) -> tuple[int, int]:
    """Create a small deck with known content. Returns (slide_count, bytes)."""
    prs = Presentation()
    blank = prs.slide_layouts[6]  # the "Blank" layout — no placeholders
    s0 = prs.slides.add_slide(blank)
    _add_textbox(s0, _SAMPLE_HEADING)
    _add_textbox(s0, "This deck contains an old figure that needs correcting.")
    _add_textbox(s0, "Status: old value", bold=True)
    s1 = prs.slides.add_slide(blank)
    _add_textbox(s1, "The old total will be updated in the old table below.")
    _add_textbox(s1, "old label")
    prs.core_properties.title = "Quarterly Report"
    prs.core_properties.author = "manual_presentation_test"
    prs.core_properties.subject = "M7 sample"
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    slide_count = len(Presentation(str(path)).slides)
    return slide_count, path.stat().st_size


# Maps a typed command to (action_type, (target, params)).
def _build(cmd: str, args: list[str], resolve):
    if cmd == "slides":
        return "presentation.slide_count", (str(resolve(args[0])), {})
    if cmd == "read":
        params = {"max_chars": int(args[1])} if len(args) > 1 else {}
        return "presentation.read_text", (str(resolve(args[0])), params)
    if cmd == "meta":
        return "presentation.get_metadata", (str(resolve(args[0])), {})
    if cmd == "find":
        if len(args) < 2:
            raise ValueError("find needs <path> <query>")
        return "presentation.find", (str(resolve(args[0])), {"query": args[1]})
    if cmd == "replace":
        if len(args) < 3:
            raise ValueError("replace needs <path> <find> <replace> [count]")
        params = {"find": args[1], "replace": args[2]}
        if len(args) > 3:
            params["count"] = int(args[3])
        return "presentation.replace_text", (str(resolve(args[0])), params)
    if cmd == "replaceas":
        if len(args) < 4:
            raise ValueError("replaceas needs <path> <find> <replace> <out> [count]")
        params = {"find": args[1], "replace": args[2], "save_as": str(resolve(args[3]))}
        if len(args) > 4:
            params["count"] = int(args[4])
        return "presentation.replace_text", (str(resolve(args[0])), params)
    raise ValueError(f"Unknown command: {cmd!r} (type 'help')")


def _print_result(result, sink: InMemoryAuditSink) -> None:
    print(f"\n  status      : {result.status.value}")
    if result.evidence:
        print(f"  evidence    : {result.evidence}")
    if result.verification is not None:
        v = result.verification
        print(f"  verification: {v.status.value} ({v.method})")
        if v.message:
            print(f"                {v.message}")
        if v.expected is not None or v.observed is not None:
            print(f"                expected={v.expected!r} observed={v.observed!r}")
    if result.error is not None:
        print(f"  error       : [{result.error.code}] {result.error.message}")
    print("  audit trail :")
    for event in sink.events:
        tag = f" ({event.outcome})" if event.outcome else ""
        print(f"      - {event.event_type.value}{tag}: {event.summary}")
    print()


async def _run(sandbox: Path) -> None:
    registry = ActionRegistry()
    register_presentation_executor(registry)
    verification = VerificationRegistry()
    register_presentation_verifiers(verification)
    audit = InMemoryAuditSink()
    dispatcher = Dispatcher(registry, AllowAllPolicy(), verification=verification, audit=audit)

    task_id = uuid.uuid4().hex[:8]
    seq = 0
    resolve = lambda p: (Path(p) if Path(p).is_absolute() else sandbox / p)

    print(f"Sandbox: {sandbox}")
    print("Manual presentation.* tester. Type 'help' for commands, 'sample' to make a test deck, 'quit' to exit.")

    while True:
        try:
            line = input("pptx> ").strip()
        except (EOFError, KeyboardInterrupt):
            print()
            break
        if not line:
            continue
        if line in ("quit", "exit"):
            break
        if line == "help":
            print(HELP)
            continue

        try:
            parts = shlex.split(line)
        except ValueError as exc:  # e.g. unbalanced quotes
            print(f"  parse error: {exc}")
            continue
        cmd, args = parts[0], parts[1:]

        # Local (non-pipeline) conveniences.
        if cmd == "sample":
            name = args[0] if args else "sample.pptx"
            target = resolve(name)
            count, size = _make_sample(target)
            print(f"  created {target} (slides={count}, {size} bytes)")
            continue
        if cmd == "ls":
            decks = sorted(p.name for p in sandbox.glob("*.pptx"))
            print(f"  {sandbox}:")
            print("    " + ("  ".join(decks) if decks else "(no decks — run 'sample')"))
            continue

        try:
            action_type, (target, params) = _build(cmd, args, resolve)
        except (ValueError, IndexError) as exc:
            print(f"  {exc}")
            continue

        action = Action(
            action_id=uuid.uuid4().hex[:8],
            task_id=task_id,
            sequence=seq,
            type=action_type,
            target=target,
            parameters=params,
            reason=f"manual: {line}",
        )
        seq += 1

        audit.events.clear()  # show only this action's events
        result = await dispatcher.dispatch(action)
        _print_result(result, audit)


def main() -> None:
    parser = argparse.ArgumentParser(description="Interactive manual tester for presentation.* actions")
    parser.add_argument(
        "--sandbox",
        default=str(Path.cwd() / "sandbox"),
        help="Directory to operate in (created if missing). Default: ./sandbox",
    )
    ns = parser.parse_args()
    sandbox = Path(ns.sandbox).resolve()
    sandbox.mkdir(parents=True, exist_ok=True)
    asyncio.run(_run(sandbox))


if __name__ == "__main__":
    main()
