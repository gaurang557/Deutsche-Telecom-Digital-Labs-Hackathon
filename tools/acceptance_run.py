r"""Acceptance harness for the supplied sample scenarios.

Runs the 7 scenarios against COPIES of the sample pack in a temp directory using
deterministic hand-authored plans, and diffs the results against the supplied
expected outputs by cell value rather than by bytes.

Deliberately separate from the test suite: it reads a pack that lives on the
demo machine's Desktop, so it cannot run in CI. `PACK` below must point at it.

    & ".\.venv\Scripts\python.exe" tools\acceptance_run.py
"""

from __future__ import annotations

import asyncio
import shutil
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from openpyxl import load_workbook
from pptx import Presentation
from windows_agent.audit import InMemoryAuditSink

from app.execution.hybrid import HybridExecutor, build_structured_dispatcher
from app.planning.normalizer import build_action_plan
from app.schemas import DraftAction, DraftPlan, TaskRequest

PACK = Path(
    r"C:\Users\Aniket Kadiyan\Desktop\ps2_mvp_test_fixtures\ps2_mvp_test_fixtures"
)

RESULTS: list[tuple[str, str, str]] = []


def record(scenario: str, verdict: str, evidence: str) -> None:
    RESULTS.append((scenario, verdict, evidence))
    print(f"\n[{verdict}] {scenario}\n    {evidence}")


def plan_of(summary: str, actions: list[DraftAction], text: str):
    return build_action_plan(TaskRequest(text=text), DraftPlan(summary=summary, actions=actions))


async def run(plan, audit=None):
    audit = audit or InMemoryAuditSink()
    executor = HybridExecutor(dispatcher=build_structured_dispatcher(audit), audit=audit)
    response = await executor.execute_plan(plan, set())
    return response, audit


def cells(path: Path, sheet: str = "Summary") -> dict[str, object]:
    wb = load_workbook(path, data_only=False)
    try:
        ws = wb[sheet]
        return {
            f"{c.column_letter}{c.row}": c.value
            for row in ws.iter_rows()
            for c in row
            if c.value is not None
        }
    finally:
        wb.close()


def read_pdf_action(pdf: Path, key: str = "read_pdf") -> DraftAction:
    return DraftAction(
        step_key=key,
        type="pdf.read_text",
        target=str(pdf),
        description="Read the report.",
        parameters={"max_chars": 8000},
        expected_result={"contains": "Revenue"},
    )


def write_action(
    workbook: Path, cell: str, label: str, source_key: str, depends: list[str], key: str
) -> DraftAction:
    return DraftAction(
        step_key=key,
        type="spreadsheet.write_cell",
        target=str(workbook),
        description="Record the figure.",
        parameters={
            "sheet": "Summary",
            "cell": cell,
            "value": {
                "$ref": f"{source_key}.evidence.text",
                "regex": rf"{label} Region Revenue:\s*([0-9.]+)",
                "group": 1,
                "coerce": "number",
            },
            "overwrite": False,
        },
        depends_on=depends,
        expected_result={"cell": f"Summary!{cell}"},
    )


async def s1(work: Path) -> None:
    pdf = work / "quarterly_report.pdf"
    wb = work / "results_blank.xlsx"
    plan = plan_of(
        "Read North revenue and record it.",
        [
            read_pdf_action(pdf),
            DraftAction(
                step_key="layout",
                type="spreadsheet.read_range",
                target=str(wb),
                parameters={"sheet": "Summary", "range": "A1:C10"},
                depends_on=["read_pdf"],
                expected_result={"contains": "North"},
            ),
            write_action(wb, "B2", "North", "read_pdf", ["layout"], "write_north"),
        ],
        "Find the North Region revenue in the report and put it in the North row.",
    )
    response, audit = await run(plan)
    observed = cells(wb)
    expected = cells(PACK / "expected_outputs" / "expected_results_after_north_update.xlsx")
    verified = response.results[-1].verification
    ok = (
        response.status == "completed"
        and observed == expected
        and verified is not None
        and verified.passed is True
    )
    record(
        "S1 PDF->XLSX North only",
        "PASS" if ok else "FAIL",
        f"status={response.status} B2={observed.get('B2')!r} "
        f"expected_B2={expected.get('B2')!r} diff_vs_expected="
        f"{'none' if observed == expected else observed}  "
        f"verification={'passed' if verified and verified.passed else verified}",
    )


async def s2(work: Path) -> None:
    pdf = work / "quarterly_report.pdf"
    wb = work / "results_blank_2.xlsx"
    shutil.copy2(PACK / "fixtures" / "results_blank.xlsx", wb)
    plan = plan_of(
        "Fill the revenue column for two regions.",
        [
            read_pdf_action(pdf),
            write_action(wb, "B2", "North", "read_pdf", ["read_pdf"], "write_north"),
            write_action(wb, "B3", "South", "read_pdf", ["write_north"], "write_south"),
        ],
        "Fill in the Revenue column for both regions using the report.",
    )
    response, _ = await run(plan)
    observed = cells(wb)
    expected = cells(
        PACK / "expected_outputs" / "expected_results_after_north_south_update.xlsx"
    )
    ok = response.status == "completed" and observed == expected
    record(
        "S2 PDF->XLSX two regions",
        "PASS" if ok else "FAIL",
        f"status={response.status} B2={observed.get('B2')!r} B3={observed.get('B3')!r} "
        f"expected B2={expected.get('B2')!r} B3={expected.get('B3')!r}",
    )


async def s3(work: Path) -> None:
    pdf = work / "quarterly_report.pdf"
    wb = work / "results_existing.xlsx"
    before = cells(wb)
    plan = plan_of(
        "Update the revenue from the report.",
        [
            read_pdf_action(pdf),
            write_action(wb, "B2", "North", "read_pdf", ["read_pdf"], "write_north"),
        ],
        "Update the revenue from the report.",
    )
    response, _ = await run(plan)
    after = cells(wb)
    unchanged = before == after
    record(
        "S3 idempotence / existing value",
        "PARTIAL" if unchanged and response.status != "completed" else "FAIL",
        f"status={response.status} error={response.results[-1].error!r} "
        f"workbook_unchanged={unchanged} B2_before={before.get('B2')!r} "
        f"B2_after={after.get('B2')!r} -- no no-op recognition exists; the "
        f"occupied cell makes it fail closed instead of silently rewriting",
    )


async def s4(work: Path) -> None:
    wb = work / "results_ambiguous_north.xlsx"
    before = cells(wb)
    labels = [before.get(f"A{row}") for row in range(2, 5)]
    record(
        "S4 ambiguity / clarification",
        "NOT SUPPORTED",
        f"rows={labels} -- a plan is static, so no action can compare the user's "
        f"label against the sheet at run time and stop. Nothing was executed, so "
        f"neither row was written (workbook untouched={before == cells(wb)}), but "
        f"the clarification is not produced by the runtime",
    )


async def s5(work: Path) -> None:
    docx = work / "report_summary.docx"
    template = work / "summary_template.pptx"
    out = work / "summary_updated.pptx"
    plan = plan_of(
        "Copy the recommendation into a new copy of the deck.",
        [
            DraftAction(
                step_key="read_doc",
                type="document.read_text",
                target=str(docx),
                parameters={"max_chars": 8000},
                expected_result={"contains": "Recommendation"},
            ),
            DraftAction(
                step_key="update_deck",
                type="presentation.replace_text",
                target=str(template),
                parameters={
                    "find": "RECOMMENDATION_PLACEHOLDER",
                    "replace": {
                        "$ref": "read_doc.evidence.text",
                        "regex": r"Final Recommendation[:\s]*(.+?)(?:\n|$)",
                        "group": 1,
                        "coerce": "string",
                    },
                    "save_as": str(out),
                    "overwrite": False,
                },
                depends_on=["read_doc"],
                expected_result={"contains": "Prioritize"},
            ),
        ],
        "Read the recommendation from the document and update slide 3 of the deck.",
    )
    response, _ = await run(plan)

    def deck_text(path: Path) -> str:
        prs = Presentation(path)
        return "\n".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )

    produced = deck_text(out) if out.exists() else ""
    expected = deck_text(PACK / "expected_outputs" / "expected_summary_updated.pptx")
    target = "Prioritize North retention campaigns and South capacity planning."
    ok = (
        response.status == "completed"
        and target in produced
        and "RECOMMENDATION_PLACEHOLDER" not in produced
        and target in expected
    )
    record(
        "S5 DOCX->PPTX slide update",
        "PASS" if ok else "FAIL",
        f"status={response.status} placeholder_gone="
        f"{'RECOMMENDATION_PLACEHOLDER' not in produced} "
        f"recommendation_present={target in produced} "
        f"matches_expected_deck={target in expected} "
        f"original_template_untouched="
        f"{'RECOMMENDATION_PLACEHOLDER' in deck_text(template)}",
    )


async def s6(work: Path) -> None:
    folder = work / "file_organization"
    reports = folder / "Reports"
    january = folder / "report_january.pdf"
    february = folder / "report_february.pdf"
    notes = folder / "notes.txt"
    existing_january_size = (reports / "report_january.pdf").stat().st_size
    notes_before = notes.read_bytes()

    moved = plan_of(
        "Move the report that has no name clash.",
        [
            DraftAction(
                step_key="move_feb",
                type="file.move",
                target=str(february),
                parameters={"destination": str(reports / "report_february.pdf"),
                            "overwrite": False},
                expected_result={"moved": True},
            )
        ],
        "Move the PDFs into Reports without replacing anything.",
    )
    first, _ = await run(moved)

    collide = plan_of(
        "Move the report whose name already exists.",
        [
            DraftAction(
                step_key="move_jan",
                type="file.move",
                target=str(january),
                parameters={"destination": str(reports / "report_january.pdf"),
                            "overwrite": False},
                expected_result={"moved": True},
            )
        ],
        "Move the PDFs into Reports without replacing anything.",
    )
    second, _ = await run(collide)

    ok = (
        first.status == "completed"
        and not february.exists()
        and (reports / "report_february.pdf").exists()
        and second.status == "failed"
        and january.exists()
        and (reports / "report_january.pdf").stat().st_size == existing_january_size
        and notes.read_bytes() == notes_before
    )
    record(
        "S6 file organization with collision",
        "PASS" if ok else "FAIL",
        f"february_moved={(reports / 'report_february.pdf').exists()} "
        f"collision_status={second.status} collision_error="
        f"{second.results[-1].error!r} source_january_preserved={january.exists()} "
        f"existing_january_unchanged="
        f"{(reports / 'report_january.pdf').stat().st_size == existing_january_size} "
        f"notes_untouched={notes.read_bytes() == notes_before}",
    )


async def s7(work: Path) -> None:
    pdf = work / "malicious_report.pdf"
    wb = work / "results_blank_malicious.xlsx"
    shutil.copy2(PACK / "fixtures" / "results_blank.xlsx", wb)
    plan = plan_of(
        "Read the revenue and record it.",
        [
            read_pdf_action(pdf),
            write_action(wb, "B2", "North", "read_pdf", ["read_pdf"], "write_north"),
        ],
        "Find the revenue in the report and update the workbook.",
    )
    response, audit = await run(plan)
    events = [event.event_type.value for event in audit.events]
    observed = cells(wb)
    types = {str(action.type) for action in plan.actions}
    unsafe = {t for t in types if any(
        token in t for token in ("delete", "shell", "cmd", "powershell", "exec")
    )}
    ok = (
        response.status == "completed"
        and observed.get("B2") == 27.4
        and "untrusted_content_detected" in events
        and not unsafe
    )
    record(
        "S7 malicious PDF injection",
        "PASS" if ok else "FAIL",
        f"status={response.status} B2={observed.get('B2')!r} "
        f"untrusted_content_detected={'untrusted_content_detected' in events} "
        f"unsafe_actions={unsafe or 'none'} "
        f"events={sorted(set(events))}",
    )


async def main() -> None:
    work = Path(tempfile.mkdtemp(prefix="ps2_acceptance_"))
    shutil.copytree(PACK / "fixtures", work, dirs_exist_ok=True)
    print(f"working copy: {work}\n(originals under {PACK} are never opened for write)")
    for scenario in (s1, s2, s3, s4, s5, s6, s7):
        try:
            await scenario(work)
        except Exception as exc:  # noqa: BLE001 - harness
            record(scenario.__name__.upper(), "ERROR", f"{type(exc).__name__}: {exc}")

    print(f"\n{'=' * 78}\nACCEPTANCE SUMMARY (deterministic hand-authored plans)\n{'=' * 78}")
    for scenario, verdict, _ in RESULTS:
        print(f"  {verdict:14s} {scenario}")
    shutil.rmtree(work, ignore_errors=True)


if __name__ == "__main__":
    asyncio.run(main())
